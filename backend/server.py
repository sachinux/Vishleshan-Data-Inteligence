from fastapi import FastAPI, APIRouter, UploadFile, File, Form, HTTPException, BackgroundTasks
from fastapi.responses import StreamingResponse
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
from pathlib import Path
from pydantic import BaseModel, Field, ConfigDict
from typing import List, Optional, Dict, Any
import uuid
from datetime import datetime, timezone
import pandas as pd
import numpy as np
import json
import io
import re
import httpx
import pdfplumber
import openpyxl
from emergentintegrations.llm.chat import LlmChat, UserMessage
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch
from pptx import Presentation
from pptx.util import Inches, Pt

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

# MongoDB connection
mongo_url = os.environ['MONGO_URL']
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ['DB_NAME']]

# LLM API Key
LLM_API_KEY = os.environ.get('EMERGENT_LLM_KEY', '')

# Create the main app
app = FastAPI()
api_router = APIRouter(prefix="/api")

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# In-memory storage for datasets (in production, use Redis or file storage)
datasets_store: Dict[str, pd.DataFrame] = {}

# ============== MODELS ==============

class Workspace(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    name: str
    description: Optional[str] = ""
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    updated_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

class WorkspaceCreate(BaseModel):
    name: str
    description: Optional[str] = ""

class DatasetFile(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    workspace_id: str
    filename: str
    file_type: str
    row_count: int = 0
    column_count: int = 0
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

class ColumnProfile(BaseModel):
    name: str
    dtype: str
    null_count: int
    null_percentage: float
    unique_count: int
    min_value: Optional[str] = None
    max_value: Optional[str] = None
    mean_value: Optional[float] = None
    sample_values: List[str] = []

class DataProfile(BaseModel):
    dataset_id: str
    row_count: int
    column_count: int
    columns: List[ColumnProfile]
    memory_usage: str

class ChatMessage(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    workspace_id: str
    role: str  # user, assistant
    content: str
    plan: Optional[str] = None
    code: Optional[str] = None
    table_data: Optional[Dict[str, Any]] = None
    chart_config: Optional[Dict[str, Any]] = None
    error: Optional[str] = None
    suggestions: List[str] = []
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    # 3-Layer System
    analysis_success: bool = True  # Whether analysis succeeded
    analysis_method: str = "auto"  # auto, statistical, aggregation, chart_only
    layer1_insight: Optional[Dict[str, Any]] = None  # Business Intelligence layer (user visible)
    layer2_reasoning: Optional[Dict[str, Any]] = None  # AI Reasoning layer (collapsible)
    layer3_runtime: Optional[Dict[str, Any]] = None  # Runtime execution layer (collapsible/hidden)
    confidence_score: Optional[int] = None  # Only set when analysis succeeds
    alternative_methods: List[str] = []  # Available methods to try if failed
    # Model Orchestrator - AI selects best analysis method
    model_selection: Optional[Dict[str, Any]] = None  # {selected, reason, alternatives: [{name, score, description}]}

class ChatRequest(BaseModel):
    workspace_id: str
    message: str
    dataset_id: Optional[str] = None
    context: Optional[str] = None  # Custom AI instructions
    response_style: Optional[str] = None  # Response style preference

class ChatSettings(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    workspace_id: str
    context: str = ""  # Up to 1000 chars
    response_style: str = ""  # Up to 50 chars
    updated_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

class ChatSettingsUpdate(BaseModel):
    context: Optional[str] = None
    response_style: Optional[str] = None

class StoryTile(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    workspace_id: str
    title: str
    key_metrics: List[str] = []
    explanation: str
    chart_config: Optional[Dict[str, Any]] = None
    table_data: Optional[Dict[str, Any]] = None
    tags: List[str] = []
    source_message_id: Optional[str] = None
    # Enhanced fields for actionable insights
    action_items: List[Dict[str, Any]] = []  # {text, priority: HIGH/MEDIUM/LOW, category}
    impact_score: str = "MEDIUM"  # HIGH, MEDIUM, LOW
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

class StoryTileCreate(BaseModel):
    workspace_id: str
    title: str
    key_metrics: List[str] = []
    explanation: str
    chart_config: Optional[Dict[str, Any]] = None
    table_data: Optional[Dict[str, Any]] = None
    tags: List[str] = []
    source_message_id: Optional[str] = None
    action_items: List[Dict[str, Any]] = []
    impact_score: str = "MEDIUM"

class StoryboardFrame(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    title: str
    summary: str
    tile_refs: List[str] = []
    narrative_notes: str = ""
    order: int = 0
    # Enhanced fields
    action_items: List[Dict[str, Any]] = []  # {text, priority, completed, category}
    kpis: List[Dict[str, Any]] = []  # {label, value, status: green/yellow/red, trend}

class KPI(BaseModel):
    label: str
    value: str
    status: str = "green"  # green, yellow, red
    trend: str = ""  # up, down, stable
    description: str = ""

class ActionItem(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    text: str
    priority: str = "MEDIUM"  # HIGH, MEDIUM, LOW
    category: str = ""  # e.g., "Marketing", "Operations", "Finance"
    completed: bool = False
    due_date: Optional[str] = None

class StakeholderView(BaseModel):
    type: str  # executive, manager, analyst
    summary: str
    key_points: List[str] = []
    recommended_actions: List[str] = []

class Storyboard(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    workspace_id: str
    title: str
    frames: List[StoryboardFrame] = []
    # Enhanced fields for actionable storyboards
    executive_summary: str = ""
    kpis: List[Dict[str, Any]] = []  # Dashboard KPIs
    action_items: List[Dict[str, Any]] = []  # Master action item list
    stakeholder_views: Dict[str, Dict[str, Any]] = {}  # {executive: {...}, manager: {...}, analyst: {...}}
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    updated_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

class StoryboardCreate(BaseModel):
    workspace_id: str
    title: str

class StoryboardUpdate(BaseModel):
    title: Optional[str] = None
    frames: Optional[List[StoryboardFrame]] = None

# ============== HELPER FUNCTIONS ==============

def serialize_datetime(obj):
    if isinstance(obj, datetime):
        return obj.isoformat()
    return obj

def prepare_for_mongo(doc: dict) -> dict:
    """Prepare a document for MongoDB storage"""
    result = {}
    for key, value in doc.items():
        if isinstance(value, datetime):
            result[key] = value.isoformat()
        elif isinstance(value, list):
            result[key] = [prepare_for_mongo(v) if isinstance(v, dict) else serialize_datetime(v) for v in value]
        elif isinstance(value, dict):
            result[key] = prepare_for_mongo(value)
        else:
            result[key] = value
    return result

def parse_datetime_fields(doc: dict, fields: List[str]) -> dict:
    """Parse datetime fields from ISO strings"""
    for field in fields:
        if field in doc and isinstance(doc[field], str):
            doc[field] = datetime.fromisoformat(doc[field])
    return doc

async def get_llm_response(prompt: str, system_message: str = "You are a data analysis assistant.") -> str:
    """Get response from LLM"""
    if not LLM_API_KEY:
        return "LLM API key not configured"
    
    chat = LlmChat(
        api_key=LLM_API_KEY,
        session_id=str(uuid.uuid4()),
        system_message=system_message
    ).with_model("openai", "gpt-5.2")
    
    user_message = UserMessage(text=prompt)
    response = await chat.send_message(user_message)
    return response

def profile_dataframe(df: pd.DataFrame) -> DataProfile:
    """Profile a pandas DataFrame"""
    columns = []
    for col in df.columns:
        col_data = df[col]
        null_count = int(col_data.isnull().sum())
        null_pct = (null_count / len(df)) * 100 if len(df) > 0 else 0
        unique_count = int(col_data.nunique())
        
        col_profile = ColumnProfile(
            name=str(col),
            dtype=str(col_data.dtype),
            null_count=null_count,
            null_percentage=round(null_pct, 2),
            unique_count=unique_count,
            sample_values=[str(v) for v in col_data.dropna().head(5).tolist()]
        )
        
        # Add numeric statistics
        if pd.api.types.is_numeric_dtype(col_data):
            col_profile.min_value = str(col_data.min()) if not col_data.isnull().all() else None
            col_profile.max_value = str(col_data.max()) if not col_data.isnull().all() else None
            col_profile.mean_value = float(col_data.mean()) if not col_data.isnull().all() else None
        else:
            col_profile.min_value = str(col_data.min()) if not col_data.isnull().all() else None
            col_profile.max_value = str(col_data.max()) if not col_data.isnull().all() else None
        
        columns.append(col_profile)
    
    memory = df.memory_usage(deep=True).sum()
    if memory > 1024 * 1024:
        memory_str = f"{memory / (1024 * 1024):.2f} MB"
    else:
        memory_str = f"{memory / 1024:.2f} KB"
    
    return DataProfile(
        dataset_id="",
        row_count=len(df),
        column_count=len(df.columns),
        columns=columns,
        memory_usage=memory_str
    )

def parse_google_sheets_url(url: str) -> Optional[str]:
    """Extract the spreadsheet ID from a Google Sheets URL"""
    patterns = [
        r'/spreadsheets/d/([a-zA-Z0-9-_]+)',
        r'id=([a-zA-Z0-9-_]+)'
    ]
    for pattern in patterns:
        match = re.search(pattern, url)
        if match:
            return match.group(1)
    return None

async def fetch_google_sheet(url: str) -> pd.DataFrame:
    """Fetch data from a public Google Sheet"""
    sheet_id = parse_google_sheets_url(url)
    if not sheet_id:
        raise ValueError("Invalid Google Sheets URL")
    
    export_url = f"https://docs.google.com/spreadsheets/d/{sheet_id}/export?format=csv"
    
    async with httpx.AsyncClient(follow_redirects=True) as client_http:
        response = await client_http.get(export_url)
        if response.status_code != 200:
            raise ValueError("Could not fetch Google Sheet. Make sure it's publicly accessible.")
        
        df = pd.read_csv(io.StringIO(response.text))
        return df

def extract_pdf_content(file_bytes: bytes) -> tuple[str, List[Dict]]:
    """Extract text and tables from PDF"""
    text_content = []
    tables = []
    
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page_num, page in enumerate(pdf.pages):
            # Extract text
            text = page.extract_text()
            if text:
                text_content.append(f"[Page {page_num + 1}]\n{text}")
            
            # Extract tables
            page_tables = page.extract_tables()
            for table in page_tables:
                if table and len(table) > 1:
                    headers = table[0]
                    data = table[1:]
                    tables.append({
                        "page": page_num + 1,
                        "headers": headers,
                        "data": data
                    })
    
    return "\n\n".join(text_content), tables

def execute_data_query(df: pd.DataFrame, query_code: str) -> Dict[str, Any]:
    """Execute a data query on a DataFrame safely"""
    try:
        # Allowed modules for import
        ALLOWED_MODULES = {
            'pandas': pd,
            'pd': pd,
            'numpy': np,
            'np': np,
            'math': __import__('math'),
            'statistics': __import__('statistics'),
            'datetime': __import__('datetime'),
            'collections': __import__('collections'),
            're': __import__('re'),
        }
        
        def safe_import(name, globals=None, locals=None, fromlist=(), level=0):
            """Controlled import that only allows whitelisted modules"""
            if name in ALLOWED_MODULES:
                return ALLOWED_MODULES[name]
            # Handle 'from X import Y' style imports
            base_module = name.split('.')[0]
            if base_module in ALLOWED_MODULES:
                return ALLOWED_MODULES[base_module]
            raise ImportError(f"Import of '{name}' is not allowed in this environment")
        
        # Create a safe execution environment with necessary builtins
        safe_builtins = {
            '__import__': safe_import,  # Controlled import function
            'len': len,
            'range': range,
            'enumerate': enumerate,
            'zip': zip,
            'map': map,
            'filter': filter,
            'sorted': sorted,
            'reversed': reversed,
            'sum': sum,
            'min': min,
            'max': max,
            'abs': abs,
            'round': round,
            'int': int,
            'float': float,
            'str': str,
            'bool': bool,
            'list': list,
            'dict': dict,
            'tuple': tuple,
            'set': set,
            'type': type,
            'isinstance': isinstance,
            'hasattr': hasattr,
            'getattr': getattr,
            'setattr': setattr,
            'print': print,
            'any': any,
            'all': all,
            'slice': slice,
            'repr': repr,
            'format': format,
            'ord': ord,
            'chr': chr,
            'divmod': divmod,
            'pow': pow,
            'True': True,
            'False': False,
            'None': None,
        }
        
        # Pre-import commonly used modules into local namespace
        local_vars = {
            "df": df, 
            "pd": pd, 
            "np": np,
            "math": __import__('math'),
            "datetime": __import__('datetime'),
            "statistics": __import__('statistics'),
        }
        
        exec(query_code, {"__builtins__": safe_builtins}, local_vars)
        
        result = local_vars.get("result", None)
        
        if result is None:
            return {"error": "No result variable found in query"}
        
        if isinstance(result, pd.DataFrame):
            return {
                "type": "dataframe",
                "data": result.head(100).to_dict(orient="records"),
                "columns": list(result.columns),
                "row_count": len(result)
            }
        elif isinstance(result, pd.Series):
            return {
                "type": "series",
                "data": result.head(100).to_dict(),
                "name": result.name
            }
        elif isinstance(result, (int, float, str, bool)):
            return {
                "type": "scalar",
                "data": result
            }
        elif isinstance(result, (np.integer, np.floating)):
            # Handle numpy scalar types
            return {
                "type": "scalar",
                "data": float(result) if isinstance(result, np.floating) else int(result)
            }
        elif isinstance(result, dict):
            return {
                "type": "dict",
                "data": result
            }
        else:
            return {
                "type": "unknown",
                "data": str(result)
            }
    except Exception as e:
        return {"error": str(e)}

# ============== API ROUTES ==============

@api_router.get("/")
async def root():
    return {"message": "Data Storyteller Studio API"}

# Workspace endpoints
@api_router.post("/workspaces", response_model=Workspace)
async def create_workspace(input: WorkspaceCreate):
    workspace = Workspace(**input.model_dump())
    doc = prepare_for_mongo(workspace.model_dump())
    await db.workspaces.insert_one(doc)
    return workspace

@api_router.get("/workspaces", response_model=List[Workspace])
async def get_workspaces():
    workspaces = await db.workspaces.find({}, {"_id": 0}).to_list(100)
    for ws in workspaces:
        parse_datetime_fields(ws, ["created_at", "updated_at"])
    return workspaces

@api_router.get("/workspaces/{workspace_id}", response_model=Workspace)
async def get_workspace(workspace_id: str):
    workspace = await db.workspaces.find_one({"id": workspace_id}, {"_id": 0})
    if not workspace:
        raise HTTPException(status_code=404, detail="Workspace not found")
    parse_datetime_fields(workspace, ["created_at", "updated_at"])
    return workspace

@api_router.delete("/workspaces/{workspace_id}")
async def delete_workspace(workspace_id: str):
    result = await db.workspaces.delete_one({"id": workspace_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Workspace not found")
    # Also delete related data
    await db.datasets.delete_many({"workspace_id": workspace_id})
    await db.chat_messages.delete_many({"workspace_id": workspace_id})
    await db.story_tiles.delete_many({"workspace_id": workspace_id})
    await db.storyboards.delete_many({"workspace_id": workspace_id})
    # Remove datasets from memory
    dataset_ids = [d["id"] for d in await db.datasets.find({"workspace_id": workspace_id}, {"id": 1, "_id": 0}).to_list(100)]
    for did in dataset_ids:
        if did in datasets_store:
            del datasets_store[did]
    return {"message": "Workspace deleted"}

class WorkspaceUpdate(BaseModel):
    name: Optional[str] = None
    description: Optional[str] = None

@api_router.put("/workspaces/{workspace_id}", response_model=Workspace)
async def update_workspace(workspace_id: str, update: WorkspaceUpdate):
    update_data = {k: v for k, v in update.model_dump().items() if v is not None}
    if not update_data:
        raise HTTPException(status_code=400, detail="No update data provided")
    
    update_data["updated_at"] = datetime.now(timezone.utc).isoformat()
    
    result = await db.workspaces.update_one(
        {"id": workspace_id},
        {"$set": update_data}
    )
    
    if result.matched_count == 0:
        raise HTTPException(status_code=404, detail="Workspace not found")
    
    workspace = await db.workspaces.find_one({"id": workspace_id}, {"_id": 0})
    parse_datetime_fields(workspace, ["created_at", "updated_at"])
    return workspace

@api_router.delete("/datasets/{dataset_id}")
async def delete_dataset(dataset_id: str):
    # Remove from MongoDB
    result = await db.datasets.delete_one({"id": dataset_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Dataset not found")
    
    # Remove from memory
    if dataset_id in datasets_store:
        del datasets_store[dataset_id]
    
    return {"message": "Dataset deleted"}

# File upload endpoints
@api_router.post("/datasets/upload")
async def upload_dataset(
    workspace_id: str = Form(...),
    file: UploadFile = File(...)
):
    filename = file.filename or "unknown"
    file_ext = filename.split(".")[-1].lower()
    
    try:
        content = await file.read()
        
        if file_ext == "csv":
            df = pd.read_csv(io.BytesIO(content))
            file_type = "csv"
        elif file_ext in ["xlsx", "xls"]:
            df = pd.read_excel(io.BytesIO(content))
            file_type = "excel"
        elif file_ext == "pdf":
            text_content, tables = extract_pdf_content(content)
            
            # If tables found, use the first one as the main dataset
            if tables:
                first_table = tables[0]
                df = pd.DataFrame(first_table["data"], columns=first_table["headers"])
            else:
                # Create a simple text dataset
                df = pd.DataFrame({"text_content": [text_content]})
            file_type = "pdf"
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {file_ext}")
        
        dataset_file = DatasetFile(
            workspace_id=workspace_id,
            filename=filename,
            file_type=file_type,
            row_count=len(df),
            column_count=len(df.columns)
        )
        
        # Store DataFrame in memory
        datasets_store[dataset_file.id] = df
        
        # Store metadata in MongoDB
        doc = prepare_for_mongo(dataset_file.model_dump())
        await db.datasets.insert_one(doc)
        
        # Get profile
        profile = profile_dataframe(df)
        profile.dataset_id = dataset_file.id
        
        return {
            "dataset": dataset_file.model_dump(),
            "profile": profile.model_dump()
        }
    
    except Exception as e:
        logger.error(f"Error uploading file: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.post("/datasets/google-sheets")
async def import_google_sheet(workspace_id: str = Form(...), url: str = Form(...)):
    try:
        df = await fetch_google_sheet(url)
        
        dataset_file = DatasetFile(
            workspace_id=workspace_id,
            filename=f"google_sheet_{uuid.uuid4().hex[:8]}",
            file_type="google_sheets",
            row_count=len(df),
            column_count=len(df.columns)
        )
        
        datasets_store[dataset_file.id] = df
        
        doc = prepare_for_mongo(dataset_file.model_dump())
        await db.datasets.insert_one(doc)
        
        profile = profile_dataframe(df)
        profile.dataset_id = dataset_file.id
        
        return {
            "dataset": dataset_file.model_dump(),
            "profile": profile.model_dump()
        }
    
    except Exception as e:
        logger.error(f"Error importing Google Sheet: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/datasets/{workspace_id}")
async def get_datasets(workspace_id: str):
    datasets = await db.datasets.find({"workspace_id": workspace_id}, {"_id": 0}).to_list(100)
    for ds in datasets:
        parse_datetime_fields(ds, ["created_at"])
    return datasets

@api_router.get("/datasets/{dataset_id}/profile")
async def get_dataset_profile(dataset_id: str):
    if dataset_id not in datasets_store:
        raise HTTPException(status_code=404, detail="Dataset not found in memory")
    
    df = datasets_store[dataset_id]
    profile = profile_dataframe(df)
    profile.dataset_id = dataset_id
    return profile.model_dump()

@api_router.get("/datasets/{dataset_id}/preview")
async def preview_dataset(dataset_id: str, rows: int = 10):
    if dataset_id not in datasets_store:
        raise HTTPException(status_code=404, detail="Dataset not found in memory")
    
    df = datasets_store[dataset_id]
    return {
        "columns": list(df.columns),
        "data": df.head(rows).to_dict(orient="records"),
        "total_rows": len(df)
    }

@api_router.get("/datasets/{dataset_id}/rows")
async def get_dataset_rows(dataset_id: str, page: int = 1, page_size: int = 100):
    """Get paginated rows from a dataset for the data grid view"""
    if dataset_id not in datasets_store:
        raise HTTPException(status_code=404, detail="Dataset not found in memory")
    
    df = datasets_store[dataset_id]
    total_rows = len(df)
    start_idx = (page - 1) * page_size
    end_idx = start_idx + page_size
    
    # Get the slice of data with row indices
    data_slice = df.iloc[start_idx:end_idx].copy()
    data_slice.insert(0, '_row_index', range(start_idx, min(end_idx, total_rows)))
    
    return {
        "columns": list(df.columns),
        "data": data_slice.to_dict(orient="records"),
        "total_rows": total_rows,
        "page": page,
        "page_size": page_size,
        "total_pages": (total_rows + page_size - 1) // page_size
    }

class SelectedRowsRequest(BaseModel):
    workspace_id: str
    dataset_id: str
    row_indices: List[int]
    action: str  # "narrate" or "compare"

@api_router.post("/datasets/selected-rows/analyze")
async def analyze_selected_rows(request: SelectedRowsRequest):
    """Analyze selected rows - narrate story or compare data"""
    if request.dataset_id not in datasets_store:
        raise HTTPException(status_code=404, detail="Dataset not found in memory")
    
    df = datasets_store[request.dataset_id]
    
    # Get selected rows
    selected_df = df.iloc[request.row_indices]
    
    if len(selected_df) == 0:
        raise HTTPException(status_code=400, detail="No rows selected")
    
    # Convert to string representation for LLM
    selected_data = selected_df.to_dict(orient="records")
    columns = list(df.columns)
    
    # Get chat settings for this workspace
    chat_settings = await db.chat_settings.find_one({"workspace_id": request.workspace_id}, {"_id": 0})
    context_instructions = ""
    response_style_instructions = ""
    
    if chat_settings:
        if chat_settings.get("context"):
            context_instructions = f"\n\nUser's custom instructions: {chat_settings['context']}"
        if chat_settings.get("response_style"):
            response_style_instructions = f" Respond in a {chat_settings['response_style']} manner."
    
    if request.action == "narrate":
        # Generate narrative story from selected rows
        prompt = f"""
You are a data storytelling expert.{response_style_instructions} Analyze these {len(selected_df)} selected rows and create a compelling narrative story.
{context_instructions}

Columns: {columns}
Selected Data: {json.dumps(selected_data[:20])}  # Limit to first 20 for context

Generate a JSON response with:
1. "title": A catchy title for this story (5-8 words)
2. "narrative": A compelling 2-3 paragraph story explaining the key insights, patterns, and significance of this data
3. "key_points": List of 3-5 bullet point insights
4. "chart_suggestion": Recommended chart type to visualize this data ("bar", "line", "pie", "scatter", or null)
5. "chart_config": If chart suggested, provide config with "x_column", "y_column", "title"

Return ONLY valid JSON.
"""
    else:  # compare
        if len(selected_df) < 2:
            raise HTTPException(status_code=400, detail="Need at least 2 rows to compare")
        
        # Generate comparison of selected rows
        prompt = f"""
You are a data analysis expert.{response_style_instructions} Compare these {len(selected_df)} selected rows and highlight differences and similarities.
{context_instructions}

Columns: {columns}
Selected Data: {json.dumps(selected_data[:20])}  # Limit to first 20 for context

Generate a JSON response with:
1. "title": A title for this comparison (5-8 words)
2. "summary": A brief summary of the comparison (1-2 sentences)
3. "similarities": List of key similarities between the rows
4. "differences": List of key differences between the rows, with specific values mentioned
5. "insights": Additional insights or patterns noticed
6. "recommendation": What action or conclusion can be drawn from this comparison
7. "chart_suggestion": Best chart type to visualize the comparison ("bar", "scatter", or null)
8. "chart_config": If chart suggested, provide config with columns to use

Return ONLY valid JSON.
"""
    
    try:
        llm_response = await get_llm_response(prompt)
        
        # Parse LLM response
        try:
            clean_response = llm_response.strip()
            if clean_response.startswith("```"):
                clean_response = clean_response.split("```")[1]
                if clean_response.startswith("json"):
                    clean_response = clean_response[4:]
            analysis = json.loads(clean_response.strip())
        except json.JSONDecodeError:
            analysis = {
                "title": "Data Analysis",
                "narrative" if request.action == "narrate" else "summary": "Analysis of selected data rows",
                "key_points" if request.action == "narrate" else "insights": ["Data analysis completed"],
            }
        
        # Create a chat message for this analysis
        content = ""
        if request.action == "narrate":
            content = analysis.get("narrative", "")
            if analysis.get("key_points"):
                content += "\n\n**Key Points:**\n" + "\n".join(f"• {p}" for p in analysis["key_points"])
        else:
            content = analysis.get("summary", "")
            if analysis.get("similarities"):
                content += "\n\n**Similarities:**\n" + "\n".join(f"• {s}" for s in analysis["similarities"])
            if analysis.get("differences"):
                content += "\n\n**Differences:**\n" + "\n".join(f"• {d}" for d in analysis["differences"])
            if analysis.get("recommendation"):
                content += f"\n\n**Recommendation:** {analysis['recommendation']}"
        
        # Build chart config if suggested
        chart_config = None
        if analysis.get("chart_suggestion") and analysis.get("chart_config"):
            chart_config = {
                "type": analysis["chart_suggestion"],
                **analysis["chart_config"]
            }
        
        # Prepare table data from selected rows
        table_data = {
            "type": "dataframe",
            "data": selected_data,
            "columns": columns,
            "row_count": len(selected_df)
        }
        
        # Save as chat message
        assistant_chat = ChatMessage(
            workspace_id=request.workspace_id,
            role="assistant",
            content=content,
            plan=f"{'Narrated story' if request.action == 'narrate' else 'Compared'} {len(selected_df)} selected rows",
            table_data=table_data,
            chart_config=chart_config,
            suggestions=[
                "Pin this insight",
                "Explore related patterns",
                "Generate more details"
            ]
        )
        
        doc = prepare_for_mongo(assistant_chat.model_dump())
        await db.chat_messages.insert_one(doc)
        
        return {
            "analysis": analysis,
            "message": assistant_chat.model_dump(),
            "selected_rows_count": len(selected_df)
        }
        
    except Exception as e:
        logger.error(f"Error analyzing selected rows: {e}")
        raise HTTPException(status_code=500, detail=str(e))


# Model Orchestrator - AI selects best analysis method
class ModelOrchestrator:
    """AI-powered model selection based on query and data characteristics"""
    
    MODELS = {
        "deep_analysis": {
            "name": "Deep Analysis",
            "description": "LLM-powered code generation for complex queries",
            "icon": "brain",
            "best_for": ["complex calculations", "custom aggregations", "multi-step analysis"]
        },
        "statistical": {
            "name": "Statistical Summary",
            "description": "Fast statistical profiling without code execution",
            "icon": "bar-chart",
            "best_for": ["data overview", "summary statistics", "column profiling"]
        },
        "aggregation": {
            "name": "Aggregation Engine",
            "description": "Direct numeric calculations (sum, mean, count)",
            "icon": "calculator",
            "best_for": ["totals", "averages", "counts", "simple math"]
        },
        "chart_generator": {
            "name": "Chart Generator",
            "description": "Direct visualization without complex analysis",
            "icon": "pie-chart",
            "best_for": ["visualizations", "charts", "graphs", "plots"]
        },
        "pattern_detector": {
            "name": "Pattern Detector",
            "description": "Anomaly and trend detection algorithms",
            "icon": "trending-up",
            "best_for": ["anomalies", "outliers", "trends", "patterns"]
        }
    }
    
    @classmethod
    def select_model(cls, query: str, df: pd.DataFrame = None) -> Dict[str, Any]:
        """Select the best model based on query analysis"""
        query_lower = query.lower()
        
        # Score each model based on query keywords
        scores = {}
        
        # Deep Analysis keywords
        deep_keywords = ["calculate", "compute", "analyze", "correlation", "regression", 
                        "predict", "compare", "filter", "group by", "pivot", "transform"]
        scores["deep_analysis"] = sum(1 for k in deep_keywords if k in query_lower) * 15
        
        # Statistical keywords
        stat_keywords = ["summary", "describe", "profile", "overview", "statistics", 
                        "stats", "info", "columns", "schema", "types", "nulls"]
        scores["statistical"] = sum(1 for k in stat_keywords if k in query_lower) * 18
        
        # Aggregation keywords
        agg_keywords = ["sum", "total", "average", "mean", "count", "min", "max", 
                       "median", "how many", "how much"]
        scores["aggregation"] = sum(1 for k in agg_keywords if k in query_lower) * 20
        
        # Chart keywords
        chart_keywords = ["chart", "graph", "plot", "visualize", "visualization", 
                         "bar", "line", "pie", "scatter", "show me", "display"]
        scores["chart_generator"] = sum(1 for k in chart_keywords if k in query_lower) * 17
        
        # Pattern keywords
        pattern_keywords = ["trend", "pattern", "anomaly", "outlier", "unusual", 
                          "detect", "find patterns", "time series"]
        scores["pattern_detector"] = sum(1 for k in pattern_keywords if k in query_lower) * 16
        
        # Add baseline scores based on data characteristics
        if df is not None:
            numeric_cols = len(df.select_dtypes(include=[np.number]).columns)
            categorical_cols = len(df.select_dtypes(include=['object', 'category']).columns)
            
            # Boost statistical for profiling-type queries
            if numeric_cols > 0:
                scores["statistical"] += 5
                scores["aggregation"] += 5
            if categorical_cols > 0:
                scores["chart_generator"] += 5
        
        # Default boost for deep analysis (it's the most versatile)
        scores["deep_analysis"] += 10
        
        # Normalize scores to percentages
        total = sum(scores.values()) or 1
        percentages = {k: round((v / total) * 100) for k, v in scores.items()}
        
        # Select the best model
        selected = max(scores, key=scores.get)
        
        # Build alternatives list sorted by score
        alternatives = []
        for model_id, score in sorted(scores.items(), key=lambda x: x[1], reverse=True):
            model_info = cls.MODELS[model_id]
            alternatives.append({
                "id": model_id,
                "name": model_info["name"],
                "score": percentages[model_id],
                "description": model_info["description"],
                "icon": model_info["icon"]
            })
        
        # Generate reason for selection
        selected_model = cls.MODELS[selected]
        reason = f"Selected {selected_model['name']} ({percentages[selected]}% match) - best for {', '.join(selected_model['best_for'][:2])}"
        
        return {
            "selected": selected,
            "selected_name": selected_model["name"],
            "selected_score": percentages[selected],
            "reason": reason,
            "alternatives": alternatives,
            "query_type": cls._detect_query_type(query_lower)
        }
    
    @staticmethod
    def _detect_query_type(query: str) -> str:
        """Detect the type of query"""
        if any(k in query for k in ["chart", "graph", "plot", "visualize"]):
            return "visualization"
        elif any(k in query for k in ["sum", "total", "average", "count"]):
            return "quantitative"
        elif any(k in query for k in ["trend", "pattern", "anomaly"]):
            return "pattern_detection"
        elif any(k in query for k in ["summary", "describe", "overview"]):
            return "exploratory"
        else:
            return "analytical"


# Chat endpoints
class AnalysisMethods:
    """Pre-built analysis methods for fallback"""
    
    @staticmethod
    def statistical_summary(df: pd.DataFrame) -> Dict[str, Any]:
        """Generate statistical summary without code execution"""
        try:
            numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
            categorical_cols = df.select_dtypes(include=['object', 'category']).columns.tolist()
            
            summary = {
                "total_rows": len(df),
                "total_columns": len(df.columns),
                "numeric_columns": len(numeric_cols),
                "categorical_columns": len(categorical_cols),
                "memory_usage": f"{df.memory_usage(deep=True).sum() / 1024:.2f} KB"
            }
            
            # Numeric statistics
            if numeric_cols:
                stats = df[numeric_cols].describe().to_dict()
                summary["numeric_stats"] = stats
            
            # Categorical value counts
            if categorical_cols:
                cat_stats = {}
                for col in categorical_cols[:3]:  # Limit to 3 columns
                    cat_stats[col] = df[col].value_counts().head(5).to_dict()
                summary["categorical_stats"] = cat_stats
            
            return {
                "type": "statistical_summary",
                "data": summary,
                "success": True
            }
        except Exception as e:
            return {"type": "statistical_summary", "error": str(e), "success": False}
    
    @staticmethod
    def simple_aggregation(df: pd.DataFrame, column: str = None) -> Dict[str, Any]:
        """Simple aggregation without complex code"""
        try:
            numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
            
            if not numeric_cols:
                return {"type": "aggregation", "error": "No numeric columns found", "success": False}
            
            target_col = column if column in numeric_cols else numeric_cols[0]
            
            result = {
                "column": target_col,
                "sum": float(df[target_col].sum()),
                "mean": float(df[target_col].mean()),
                "median": float(df[target_col].median()),
                "min": float(df[target_col].min()),
                "max": float(df[target_col].max()),
                "std": float(df[target_col].std()),
                "count": int(df[target_col].count())
            }
            
            return {
                "type": "aggregation",
                "data": result,
                "success": True
            }
        except Exception as e:
            return {"type": "aggregation", "error": str(e), "success": False}
    
    @staticmethod
    def chart_only(df: pd.DataFrame) -> Dict[str, Any]:
        """Generate chart config without complex analysis"""
        try:
            numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
            categorical_cols = df.select_dtypes(include=['object', 'category']).columns.tolist()
            
            if numeric_cols and categorical_cols:
                # Bar chart with category vs numeric
                y_col = numeric_cols[0]
                x_col = categorical_cols[0]
                
                chart_data = df.groupby(x_col)[y_col].mean().head(10).reset_index()
                
                return {
                    "type": "chart_only",
                    "chart_config": {
                        "type": "bar",
                        "x_column": x_col,
                        "y_column": y_col,
                        "title": f"Average {y_col} by {x_col}"
                    },
                    "data": chart_data.to_dict(orient="records"),
                    "success": True
                }
            elif numeric_cols and len(numeric_cols) >= 2:
                # Scatter plot with two numeric columns
                return {
                    "type": "chart_only",
                    "chart_config": {
                        "type": "scatter",
                        "x_column": numeric_cols[0],
                        "y_column": numeric_cols[1],
                        "title": f"{numeric_cols[0]} vs {numeric_cols[1]}"
                    },
                    "data": df[[numeric_cols[0], numeric_cols[1]]].head(100).to_dict(orient="records"),
                    "success": True
                }
            else:
                return {"type": "chart_only", "error": "Not enough columns for chart", "success": False}
        except Exception as e:
            return {"type": "chart_only", "error": str(e), "success": False}


@api_router.post("/chat")
async def chat(request: ChatRequest):
    workspace_id = request.workspace_id
    user_message = request.message
    dataset_id = request.dataset_id
    
    # Get chat settings for this workspace
    chat_settings = await db.chat_settings.find_one({"workspace_id": workspace_id}, {"_id": 0})
    context_instructions = ""
    response_style_instructions = ""
    
    if chat_settings:
        if chat_settings.get("context"):
            context_instructions = f"\n\nUser's custom instructions: {chat_settings['context']}"
        if chat_settings.get("response_style"):
            response_style_instructions = f" Respond in a {chat_settings['response_style']} manner."
    
    # Also check request-level overrides
    if request.context:
        context_instructions = f"\n\nUser's custom instructions: {request.context}"
    if request.response_style:
        response_style_instructions = f" Respond in a {request.response_style} manner."
    
    # Save user message
    user_chat = ChatMessage(
        workspace_id=workspace_id,
        role="user",
        content=user_message
    )
    doc = prepare_for_mongo(user_chat.model_dump())
    await db.chat_messages.insert_one(doc)
    
    # Get dataset info if available
    df = None
    schema_info = ""
    if dataset_id and dataset_id in datasets_store:
        df = datasets_store[dataset_id]
        profile = profile_dataframe(df)
        schema_info = f"""
Dataset Schema:
- Rows: {profile.row_count}
- Columns: {profile.column_count}
- Column Details:
"""
        for col in profile.columns:
            schema_info += f"  * {col.name} ({col.dtype}): {col.unique_count} unique values, {col.null_percentage}% null"
            if col.mean_value is not None:
                schema_info += f", mean={col.mean_value:.2f}"
            schema_info += "\n"
    
    # Generate analysis plan using LLM
    analysis_prompt = f"""
You are a data analysis assistant.{response_style_instructions} The user asked: "{user_message}"
{context_instructions}
{schema_info}

Generate a JSON response with:
1. "plan": A brief explanation of what analysis you'll perform (1-2 sentences)
2. "code": Python code using pandas to analyze the data. The DataFrame is available as 'df'. Store the result in a variable called 'result'. The result should be a DataFrame, Series, or scalar value.
3. "chart_type": If visualization is appropriate, specify one of: "bar", "line", "scatter", "pie", "heatmap", or null if no chart needed
4. "chart_config": If chart_type is specified, provide config with "x_column", "y_column", "title", and optionally "color_by"
5. "suggestions": List of 3 follow-up questions the user might want to ask

Return ONLY valid JSON, no markdown formatting.
Example:
{{"plan": "I'll calculate the average sales by region", "code": "result = df.groupby('region')['sales'].mean()", "chart_type": "bar", "chart_config": {{"x_column": "region", "y_column": "sales", "title": "Average Sales by Region"}}, "suggestions": ["Show top 5 regions", "Compare year over year", "Filter by date range"]}}
"""
    
    # AI Model Orchestrator - Select best analysis method
    model_selection = ModelOrchestrator.select_model(user_message, df)
    
    # Initialize 3-layer response structure
    layer1_insight = {"summary": "", "recommendations": [], "key_findings": []}
    layer2_reasoning = {
        "methodology": "", 
        "steps": [], 
        "data_quality_notes": [],
        "model_selection": model_selection  # Include model selection in Layer 2
    }
    layer3_runtime = {"code": None, "execution_time_ms": 0, "error_details": None, "stack_trace": None}
    analysis_success = True
    confidence_score = None
    alternative_methods = ["statistical", "aggregation", "chart_only"]
    
    try:
        import time
        
        llm_response = await get_llm_response(analysis_prompt)
        
        # Parse LLM response
        try:
            # Clean up the response - remove markdown if present
            clean_response = llm_response.strip()
            if clean_response.startswith("```"):
                clean_response = clean_response.split("```")[1]
                if clean_response.startswith("json"):
                    clean_response = clean_response[4:]
            clean_response = clean_response.strip()
            
            analysis = json.loads(clean_response)
        except json.JSONDecodeError:
            analysis = {
                "plan": "I'll analyze your request",
                "code": "result = df.describe()" if df is not None else "result = 'No dataset loaded'",
                "chart_type": None,
                "chart_config": None,
                "suggestions": ["Upload a dataset first", "Ask about specific columns", "Request a summary"]
            }
        
        # Populate Layer 2 - AI Reasoning
        layer2_reasoning["methodology"] = analysis.get("plan", "")
        layer2_reasoning["steps"] = [
            "Parsed user query",
            "Analyzed dataset schema",
            "Generated analysis plan",
            "Created execution code"
        ]
        
        # Execute the analysis code
        result_data = None
        error = None
        layer3_runtime["code"] = analysis.get("code")
        
        if df is not None and "code" in analysis:
            exec_start = time.time()
            result_data = execute_data_query(df, analysis["code"])
            layer3_runtime["execution_time_ms"] = int((time.time() - exec_start) * 1000)
            
            if "error" in result_data:
                error = result_data["error"]
                layer3_runtime["error_details"] = error
                layer3_runtime["stack_trace"] = f"Code execution failed: {error}"
                result_data = None
                analysis_success = False
        
        # Calculate confidence score (only if successful)
        if analysis_success and result_data:
            confidence_score = 50
            if result_data.get("data") is not None:
                confidence_score += 20
            if analysis.get("chart_type"):
                confidence_score += 15
            if analysis.get("plan"):
                confidence_score += 10
            # Handle different result types
            result_data_value = result_data.get("data")
            if isinstance(result_data_value, (list, dict)):
                if len(result_data_value) > 0:
                    confidence_score += 5
            elif result_data_value is not None:  # Scalar value
                confidence_score += 5
            confidence_score = min(confidence_score, 100)
        else:
            confidence_score = None  # Don't show confidence on failure
        
        # Generate natural language answer
        system_message = f"You are a friendly data analyst explaining results to a business user.{response_style_instructions}"
        answer_prompt = f"""
The user asked: "{user_message}"
Analysis plan: {analysis.get('plan', 'N/A')}
Result: {json.dumps(result_data) if result_data else 'No result'}
Error: {error if error else 'None'}
{context_instructions}

Provide a clear, concise answer in 2-3 sentences. If there was an error, explain what might have gone wrong and suggest fixes.
"""
        answer = await get_llm_response(answer_prompt, system_message)
        
        # Populate Layer 1 - Business Intelligence
        layer1_insight["summary"] = answer
        if result_data and analysis_success:
            layer1_insight["key_findings"] = [
                f"Analyzed {result_data.get('row_count', 0)} rows of data",
                analysis.get("plan", "Analysis completed")
            ]
            if analysis.get("chart_type"):
                layer1_insight["recommendations"].append(f"View the {analysis.get('chart_type')} chart for visual insights")
        
        # Create assistant response with 3-layer structure
        assistant_chat = ChatMessage(
            workspace_id=workspace_id,
            role="assistant",
            content=answer,
            plan=analysis.get("plan"),
            code=analysis.get("code"),
            table_data=result_data if result_data and result_data.get("type") in ["dataframe", "series", "dict"] else None,
            chart_config={
                "type": analysis.get("chart_type"),
                **analysis.get("chart_config", {})
            } if analysis.get("chart_type") else None,
            error=error,
            suggestions=analysis.get("suggestions", []),
            # 3-Layer fields
            analysis_success=analysis_success,
            analysis_method="auto",
            layer1_insight=layer1_insight,
            layer2_reasoning=layer2_reasoning,
            layer3_runtime=layer3_runtime,
            confidence_score=confidence_score,
            alternative_methods=alternative_methods if not analysis_success else [],
            model_selection=model_selection
        )
        
        doc = prepare_for_mongo(assistant_chat.model_dump())
        await db.chat_messages.insert_one(doc)
        
        return assistant_chat.model_dump()
    
    except Exception as e:
        logger.error(f"Error in chat: {e}")
        
        # Populate error info in layers
        layer3_runtime["error_details"] = str(e)
        layer3_runtime["stack_trace"] = f"Exception: {type(e).__name__}: {str(e)}"
        layer1_insight["summary"] = "Analysis temporarily unavailable in this environment."
        
        error_chat = ChatMessage(
            workspace_id=workspace_id,
            role="assistant",
            content="I encountered an issue processing your request. You can try an alternative analysis method or retry.",
            error=str(e),
            suggestions=["Try a simpler question", "Check if dataset is loaded", "Ask about available columns"],
            analysis_success=False,
            analysis_method="auto",
            layer1_insight=layer1_insight,
            layer2_reasoning=layer2_reasoning,
            layer3_runtime=layer3_runtime,
            confidence_score=None,
            alternative_methods=["statistical", "aggregation", "chart_only"],
            model_selection=model_selection if 'model_selection' in dir() else None
        )
        doc = prepare_for_mongo(error_chat.model_dump())
        await db.chat_messages.insert_one(doc)
        return error_chat.model_dump()


class AlternativeAnalysisRequest(BaseModel):
    workspace_id: str
    dataset_id: str
    method: str  # statistical, aggregation, chart_only
    original_query: Optional[str] = None


@api_router.post("/chat/alternative")
async def chat_alternative(request: AlternativeAnalysisRequest):
    """Run analysis using an alternative method when primary fails"""
    workspace_id = request.workspace_id
    dataset_id = request.dataset_id
    method = request.method
    
    if dataset_id not in datasets_store:
        raise HTTPException(status_code=404, detail="Dataset not found in memory")
    
    df = datasets_store[dataset_id]
    
    # Initialize layers
    layer1_insight = {"summary": "", "recommendations": [], "key_findings": []}
    layer2_reasoning = {"methodology": "", "steps": [], "data_quality_notes": []}
    layer3_runtime = {"code": None, "execution_time_ms": 0, "error_details": None, "stack_trace": None}
    
    try:
        import time
        start_time = time.time()
        
        if method == "statistical":
            result = AnalysisMethods.statistical_summary(df)
            layer2_reasoning["methodology"] = "Generated statistical summary of all columns"
            layer2_reasoning["steps"] = ["Identified numeric and categorical columns", "Calculated descriptive statistics", "Compiled summary"]
            
            if result["success"]:
                layer1_insight["summary"] = f"Statistical summary: {result['data']['total_rows']} rows, {result['data']['total_columns']} columns"
                layer1_insight["key_findings"] = [
                    f"Dataset has {result['data']['numeric_columns']} numeric columns",
                    f"Dataset has {result['data']['categorical_columns']} categorical columns",
                    f"Memory usage: {result['data']['memory_usage']}"
                ]
                
        elif method == "aggregation":
            result = AnalysisMethods.simple_aggregation(df)
            layer2_reasoning["methodology"] = "Performed simple aggregation on numeric column"
            layer2_reasoning["steps"] = ["Selected first numeric column", "Calculated sum, mean, median, min, max", "Compiled results"]
            
            if result["success"]:
                data = result['data']
                layer1_insight["summary"] = f"Analysis of '{data['column']}': Mean = {data['mean']:.2f}, Sum = {data['sum']:.2f}"
                layer1_insight["key_findings"] = [
                    f"Mean: {data['mean']:.2f}",
                    f"Min: {data['min']:.2f}, Max: {data['max']:.2f}",
                    f"Count: {data['count']} records"
                ]
                
        elif method == "chart_only":
            result = AnalysisMethods.chart_only(df)
            layer2_reasoning["methodology"] = "Generated chart without complex analysis"
            layer2_reasoning["steps"] = ["Identified suitable columns", "Generated chart configuration", "Prepared visualization data"]
            
            if result["success"]:
                layer1_insight["summary"] = f"Created {result['chart_config']['type']} chart: {result['chart_config']['title']}"
                layer1_insight["recommendations"] = ["View the chart for visual patterns"]
        else:
            raise HTTPException(status_code=400, detail=f"Unknown method: {method}")
        
        layer3_runtime["execution_time_ms"] = int((time.time() - start_time) * 1000)
        
        # Calculate confidence
        confidence_score = 70 if result.get("success") else None
        
        # Prepare response content
        content = layer1_insight["summary"] if result.get("success") else f"Alternative analysis failed: {result.get('error', 'Unknown error')}"
        
        # Build table_data for display
        table_data = None
        if result.get("success") and result.get("data"):
            if method == "statistical":
                # Convert stats to displayable format
                table_data = {
                    "type": "dict",
                    "data": result["data"]
                }
            elif method == "aggregation":
                table_data = {
                    "type": "dict",
                    "data": result["data"]
                }
            elif method == "chart_only" and result.get("data"):
                table_data = {
                    "type": "dataframe",
                    "data": result["data"],
                    "row_count": len(result["data"])
                }
        
        # Create response message
        assistant_chat = ChatMessage(
            workspace_id=workspace_id,
            role="assistant",
            content=content,
            plan=layer2_reasoning["methodology"],
            table_data=table_data,
            chart_config=result.get("chart_config") if method == "chart_only" and result.get("success") else None,
            error=result.get("error") if not result.get("success") else None,
            suggestions=["Try another method", "Ask a specific question", "Upload more data"],
            analysis_success=result.get("success", False),
            analysis_method=method,
            layer1_insight=layer1_insight,
            layer2_reasoning=layer2_reasoning,
            layer3_runtime=layer3_runtime,
            confidence_score=confidence_score,
            alternative_methods=[]
        )
        
        doc = prepare_for_mongo(assistant_chat.model_dump())
        await db.chat_messages.insert_one(doc)
        
        return assistant_chat.model_dump()
        
    except Exception as e:
        logger.error(f"Error in alternative analysis: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/chat/{workspace_id}")
async def get_chat_history(workspace_id: str):
    messages = await db.chat_messages.find(
        {"workspace_id": workspace_id}, 
        {"_id": 0}
    ).sort("created_at", 1).to_list(1000)
    
    for msg in messages:
        parse_datetime_fields(msg, ["created_at"])
    
    return messages

# Chat Settings endpoints
@api_router.get("/chat-settings/{workspace_id}")
async def get_chat_settings(workspace_id: str):
    settings = await db.chat_settings.find_one({"workspace_id": workspace_id}, {"_id": 0})
    if not settings:
        # Return default settings if none exist
        return {"workspace_id": workspace_id, "context": "", "response_style": ""}
    parse_datetime_fields(settings, ["updated_at"])
    return settings

@api_router.put("/chat-settings/{workspace_id}")
async def update_chat_settings(workspace_id: str, update: ChatSettingsUpdate):
    update_data = {}
    
    if update.context is not None:
        # Limit to 1000 characters
        update_data["context"] = update.context[:1000]
    if update.response_style is not None:
        # Limit to 50 characters
        update_data["response_style"] = update.response_style[:50]
    
    if not update_data:
        raise HTTPException(status_code=400, detail="No update data provided")
    
    update_data["updated_at"] = datetime.now(timezone.utc).isoformat()
    update_data["workspace_id"] = workspace_id
    
    # Upsert: create if doesn't exist, update if it does
    await db.chat_settings.update_one(
        {"workspace_id": workspace_id},
        {"$set": update_data},
        upsert=True
    )
    
    settings = await db.chat_settings.find_one({"workspace_id": workspace_id}, {"_id": 0})
    parse_datetime_fields(settings, ["updated_at"])
    return settings

# Story Tiles endpoints
@api_router.post("/story-tiles", response_model=StoryTile)
async def create_story_tile(input: StoryTileCreate):
    tile = StoryTile(**input.model_dump())
    doc = prepare_for_mongo(tile.model_dump())
    await db.story_tiles.insert_one(doc)
    return tile

@api_router.post("/story-tiles/from-message")
async def create_tile_from_message(workspace_id: str = Form(...), message_id: str = Form(...)):
    """Create a story tile from a chat message"""
    message = await db.chat_messages.find_one({"id": message_id}, {"_id": 0})
    if not message:
        raise HTTPException(status_code=404, detail="Message not found")
    
    # Use LLM to generate tile content with action items
    tile_prompt = f"""
Based on this analysis result, create an actionable story tile with:
1. "title": A catchy, brief title (5-8 words)
2. "key_metrics": List of 2-3 key numbers/insights from the result
3. "explanation": A 1-2 sentence explanation for business users
4. "tags": List of 2-3 relevant tags
5. "impact_score": Rate the business impact as "HIGH", "MEDIUM", or "LOW"
6. "action_items": List of 1-3 actionable recommendations, each with:
   - "text": The action to take
   - "priority": "HIGH", "MEDIUM", or "LOW"
   - "category": Category like "Marketing", "Operations", "Finance", "Product", "Sales"

Analysis content: {message.get('content')}
Result data: {json.dumps(message.get('table_data')) if message.get('table_data') else 'N/A'}

Return ONLY valid JSON.
"""
    
    llm_response = await get_llm_response(tile_prompt)
    
    try:
        clean_response = llm_response.strip()
        if clean_response.startswith("```"):
            clean_response = clean_response.split("```")[1]
            if clean_response.startswith("json"):
                clean_response = clean_response[4:]
        tile_data = json.loads(clean_response.strip())
    except json.JSONDecodeError:
        tile_data = {
            "title": "Analysis Result",
            "key_metrics": [],
            "explanation": message.get("content", "")[:200],
            "tags": [],
            "impact_score": "MEDIUM",
            "action_items": []
        }
    
    tile = StoryTile(
        workspace_id=workspace_id,
        title=tile_data.get("title", "Analysis Result"),
        key_metrics=tile_data.get("key_metrics", []),
        explanation=tile_data.get("explanation", ""),
        chart_config=message.get("chart_config"),
        table_data=message.get("table_data"),
        tags=tile_data.get("tags", []),
        source_message_id=message_id,
        action_items=tile_data.get("action_items", []),
        impact_score=tile_data.get("impact_score", "MEDIUM")
    )
    
    doc = prepare_for_mongo(tile.model_dump())
    await db.story_tiles.insert_one(doc)
    
    return tile.model_dump()

@api_router.get("/story-tiles/{workspace_id}")
async def get_story_tiles(workspace_id: str):
    tiles = await db.story_tiles.find({"workspace_id": workspace_id}, {"_id": 0}).to_list(100)
    for tile in tiles:
        parse_datetime_fields(tile, ["created_at"])
    return tiles

@api_router.delete("/story-tiles/{tile_id}")
async def delete_story_tile(tile_id: str):
    result = await db.story_tiles.delete_one({"id": tile_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Story tile not found")
    return {"message": "Story tile deleted"}

# Storyboard endpoints
@api_router.post("/storyboards", response_model=Storyboard)
async def create_storyboard(input: StoryboardCreate):
    storyboard = Storyboard(**input.model_dump())
    doc = prepare_for_mongo(storyboard.model_dump())
    await db.storyboards.insert_one(doc)
    return storyboard

@api_router.post("/storyboards/generate")
async def generate_storyboard(workspace_id: str = Form(...), title: str = Form("Data Story")):
    """Auto-generate an actionable storyboard from existing story tiles"""
    tiles = await db.story_tiles.find({"workspace_id": workspace_id}, {"_id": 0}).to_list(100)
    
    if not tiles:
        raise HTTPException(status_code=400, detail="No story tiles found. Create some tiles first.")
    
    # Prepare tiles data with action items
    tiles_summary = json.dumps([{
        "id": t["id"], 
        "title": t["title"], 
        "explanation": t.get("explanation", ""),
        "key_metrics": t.get("key_metrics", []),
        "action_items": t.get("action_items", []),
        "impact_score": t.get("impact_score", "MEDIUM")
    } for t in tiles])
    
    # Use LLM to generate comprehensive storyboard
    storyboard_prompt = f"""
Create a comprehensive, actionable storyboard from these story tiles: {tiles_summary}

Generate a JSON response with:

1. "executive_summary": A 2-3 sentence executive summary highlighting the most critical findings and recommended actions.

2. "kpis": Array of 3-5 key performance indicators, each with:
   - "label": KPI name
   - "value": The metric value
   - "status": "green" (good), "yellow" (warning), or "red" (critical)
   - "trend": "up", "down", or "stable"
   - "description": Brief explanation

3. "action_items": Master list of prioritized action items (consolidate from tiles), each with:
   - "id": Unique ID
   - "text": Action description
   - "priority": "HIGH", "MEDIUM", or "LOW"
   - "category": Business area (Marketing, Operations, Finance, Product, Sales)
   - "completed": false

4. "frames": Array of story frames. Each frame should have:
   - "title": Frame title
   - "summary": 1-2 sentence summary
   - "tile_refs": List of tile IDs to include
   - "narrative_notes": Speaker notes for presenting
   - "action_items": Frame-specific actions
   - "kpis": Frame-specific KPIs if applicable

5. "stakeholder_views": Object with different perspectives:
   - "executive": {{"summary": "1-2 sentences for C-suite", "key_points": ["point1", "point2"], "recommended_actions": ["action1"]}}
   - "manager": {{"summary": "Tactical summary", "key_points": [...], "recommended_actions": [...]}}
   - "analyst": {{"summary": "Detailed findings", "key_points": [...], "recommended_actions": [...]}}

Organize frames into a logical narrative: Problem Statement -> Key Findings -> Recommendations -> Next Steps.
Return ONLY valid JSON.
"""
    
    llm_response = await get_llm_response(storyboard_prompt)
    
    try:
        clean_response = llm_response.strip()
        if clean_response.startswith("```"):
            clean_response = clean_response.split("```")[1]
            if clean_response.startswith("json"):
                clean_response = clean_response[4:]
        storyboard_data = json.loads(clean_response.strip())
    except json.JSONDecodeError:
        # Default structure
        storyboard_data = {
            "executive_summary": "Analysis of key business metrics reveals opportunities for improvement.",
            "kpis": [],
            "action_items": [],
            "frames": [{
                "title": "Overview",
                "summary": "Key findings from the analysis",
                "tile_refs": [t["id"] for t in tiles[:3]],
                "narrative_notes": "Let's walk through the main insights.",
                "action_items": [],
                "kpis": []
            }],
            "stakeholder_views": {
                "executive": {"summary": "High-level overview", "key_points": [], "recommended_actions": []},
                "manager": {"summary": "Tactical overview", "key_points": [], "recommended_actions": []},
                "analyst": {"summary": "Detailed findings", "key_points": [], "recommended_actions": []}
            }
        }
    
    frames = []
    for i, frame_data in enumerate(storyboard_data.get("frames", [])):
        # Ensure action_items are proper dictionaries
        frame_action_items = []
        for item in frame_data.get("action_items", []):
            if isinstance(item, dict):
                frame_action_items.append(item)
            elif isinstance(item, str):
                # If it's a string (ID reference), skip it or find it in master list
                pass
        
        # Ensure kpis are proper dictionaries
        frame_kpis = []
        for kpi in frame_data.get("kpis", []):
            if isinstance(kpi, dict):
                frame_kpis.append(kpi)
        
        frame = StoryboardFrame(
            title=frame_data.get("title", f"Frame {i+1}"),
            summary=frame_data.get("summary", ""),
            tile_refs=frame_data.get("tile_refs", []),
            narrative_notes=frame_data.get("narrative_notes", ""),
            order=i,
            action_items=frame_action_items,
            kpis=frame_kpis
        )
        frames.append(frame)
    
    # Ensure master action_items have proper IDs
    master_action_items = []
    for i, item in enumerate(storyboard_data.get("action_items", [])):
        if isinstance(item, dict):
            if not item.get("id"):
                item["id"] = f"action-{i+1}"
            if not item.get("completed"):
                item["completed"] = False
            master_action_items.append(item)
    
    storyboard = Storyboard(
        workspace_id=workspace_id,
        title=title,
        frames=frames,
        executive_summary=storyboard_data.get("executive_summary", ""),
        kpis=storyboard_data.get("kpis", []) if isinstance(storyboard_data.get("kpis"), list) else [],
        action_items=master_action_items,
        stakeholder_views=storyboard_data.get("stakeholder_views", {}) if isinstance(storyboard_data.get("stakeholder_views"), dict) else {}
    )
    
    doc = prepare_for_mongo(storyboard.model_dump())
    await db.storyboards.insert_one(doc)
    
    return storyboard.model_dump()

@api_router.get("/storyboards/{workspace_id}")
async def get_storyboards(workspace_id: str):
    storyboards = await db.storyboards.find({"workspace_id": workspace_id}, {"_id": 0}).to_list(100)
    for sb in storyboards:
        parse_datetime_fields(sb, ["created_at", "updated_at"])
    return storyboards

@api_router.get("/storyboards/detail/{storyboard_id}")
async def get_storyboard(storyboard_id: str):
    storyboard = await db.storyboards.find_one({"id": storyboard_id}, {"_id": 0})
    if not storyboard:
        raise HTTPException(status_code=404, detail="Storyboard not found")
    parse_datetime_fields(storyboard, ["created_at", "updated_at"])
    return storyboard

@api_router.put("/storyboards/{storyboard_id}")
async def update_storyboard(storyboard_id: str, update: StoryboardUpdate):
    update_data = {k: v for k, v in update.model_dump().items() if v is not None}
    if "frames" in update_data:
        update_data["frames"] = [f.model_dump() if hasattr(f, 'model_dump') else f for f in update_data["frames"]]
    update_data["updated_at"] = datetime.now(timezone.utc).isoformat()
    
    result = await db.storyboards.update_one(
        {"id": storyboard_id},
        {"$set": update_data}
    )
    
    if result.matched_count == 0:
        raise HTTPException(status_code=404, detail="Storyboard not found")
    
    return await get_storyboard(storyboard_id)

@api_router.delete("/storyboards/{storyboard_id}")
async def delete_storyboard(storyboard_id: str):
    result = await db.storyboards.delete_one({"id": storyboard_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Storyboard not found")
    return {"message": "Storyboard deleted"}

# Action Item toggle endpoint
class ActionItemToggle(BaseModel):
    action_id: str
    completed: bool

@api_router.put("/storyboards/{storyboard_id}/action-items")
async def toggle_action_item(storyboard_id: str, toggle: ActionItemToggle):
    """Toggle the completion status of an action item"""
    storyboard = await db.storyboards.find_one({"id": storyboard_id}, {"_id": 0})
    if not storyboard:
        raise HTTPException(status_code=404, detail="Storyboard not found")
    
    # Update action item in the master list
    action_items = storyboard.get("action_items", [])
    for item in action_items:
        if item.get("id") == toggle.action_id:
            item["completed"] = toggle.completed
            break
    
    # Also update in frames if present
    frames = storyboard.get("frames", [])
    for frame in frames:
        for item in frame.get("action_items", []):
            if item.get("id") == toggle.action_id:
                item["completed"] = toggle.completed
    
    await db.storyboards.update_one(
        {"id": storyboard_id},
        {"$set": {"action_items": action_items, "frames": frames, "updated_at": datetime.now(timezone.utc).isoformat()}}
    )
    
    return {"message": "Action item updated", "action_id": toggle.action_id, "completed": toggle.completed}

# Export endpoints
@api_router.post("/export/pdf/{storyboard_id}")
async def export_storyboard_pdf(storyboard_id: str):
    storyboard = await db.storyboards.find_one({"id": storyboard_id}, {"_id": 0})
    if not storyboard:
        raise HTTPException(status_code=404, detail="Storyboard not found")
    
    # Get tiles for the storyboard
    tile_ids = []
    for frame in storyboard.get("frames", []):
        tile_ids.extend(frame.get("tile_refs", []))
    
    tiles = {}
    if tile_ids:
        tiles_list = await db.story_tiles.find({"id": {"$in": tile_ids}}, {"_id": 0}).to_list(100)
        tiles = {t["id"]: t for t in tiles_list}
    
    # Generate PDF
    buffer = io.BytesIO()
    c = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter
    
    # Title page
    c.setFont("Helvetica-Bold", 24)
    c.drawCentredString(width/2, height - 2*inch, storyboard.get("title", "Data Story"))
    c.setFont("Helvetica", 12)
    c.drawCentredString(width/2, height - 2.5*inch, f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    c.showPage()
    
    # Frames
    for frame in storyboard.get("frames", []):
        c.setFont("Helvetica-Bold", 18)
        c.drawString(0.75*inch, height - 1*inch, frame.get("title", "Frame"))
        
        c.setFont("Helvetica", 12)
        y = height - 1.5*inch
        
        # Summary
        c.drawString(0.75*inch, y, "Summary:")
        y -= 0.25*inch
        summary_text = frame.get("summary", "")
        c.drawString(0.75*inch, y, summary_text[:80])
        y -= 0.5*inch
        
        # Tiles content
        for tile_id in frame.get("tile_refs", []):
            tile = tiles.get(tile_id, {})
            if tile:
                c.setFont("Helvetica-Bold", 12)
                c.drawString(0.75*inch, y, tile.get("title", "Insight"))
                y -= 0.25*inch
                
                c.setFont("Helvetica", 10)
                for metric in tile.get("key_metrics", [])[:3]:
                    c.drawString(1*inch, y, f"• {metric}")
                    y -= 0.2*inch
                
                explanation = tile.get("explanation", "")[:200]
                c.drawString(0.75*inch, y, explanation)
                y -= 0.4*inch
                
                if y < 2*inch:
                    break
        
        # Narrative notes
        y -= 0.25*inch
        c.setFont("Helvetica-Oblique", 10)
        c.drawString(0.75*inch, y, "Notes: " + frame.get("narrative_notes", "")[:100])
        
        c.showPage()
    
    c.save()
    buffer.seek(0)
    
    return StreamingResponse(
        buffer,
        media_type="application/pdf",
        headers={"Content-Disposition": f"attachment; filename={storyboard.get('title', 'storyboard')}.pdf"}
    )

@api_router.post("/export/pptx/{storyboard_id}")
async def export_storyboard_pptx(storyboard_id: str):
    storyboard = await db.storyboards.find_one({"id": storyboard_id}, {"_id": 0})
    if not storyboard:
        raise HTTPException(status_code=404, detail="Storyboard not found")
    
    # Get tiles
    tile_ids = []
    for frame in storyboard.get("frames", []):
        tile_ids.extend(frame.get("tile_refs", []))
    
    tiles = {}
    if tile_ids:
        tiles_list = await db.story_tiles.find({"id": {"$in": tile_ids}}, {"_id": 0}).to_list(100)
        tiles = {t["id"]: t for t in tiles_list}
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(title_slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = storyboard.get("title", "Data Story")
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    
    # Frame slides
    for frame in storyboard.get("frames", []):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Frame title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = frame.get("title", "Frame")
        p.font.size = Pt(32)
        p.font.bold = True
        
        # Summary
        summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.6))
        tf = summary_box.text_frame
        p = tf.paragraphs[0]
        p.text = frame.get("summary", "")
        p.font.size = Pt(18)
        
        # Tiles content
        y_pos = 2.0
        for tile_id in frame.get("tile_refs", [])[:2]:
            tile = tiles.get(tile_id, {})
            if tile:
                tile_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(6), Inches(2))
                tf = tile_box.text_frame
                tf.word_wrap = True
                
                p = tf.paragraphs[0]
                p.text = tile.get("title", "")
                p.font.size = Pt(20)
                p.font.bold = True
                
                for metric in tile.get("key_metrics", [])[:3]:
                    p = tf.add_paragraph()
                    p.text = f"• {metric}"
                    p.font.size = Pt(14)
                
                p = tf.add_paragraph()
                p.text = tile.get("explanation", "")[:150]
                p.font.size = Pt(12)
                
                y_pos += 2.5
        
        # Narrative notes at bottom
        notes_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.8))
        tf = notes_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Notes: {frame.get('narrative_notes', '')}"
        p.font.size = Pt(12)
        p.font.italic = True
    
    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={storyboard.get('title', 'storyboard')}.pptx"}
    )

@api_router.get("/export/json/{storyboard_id}")
async def export_storyboard_json(storyboard_id: str):
    storyboard = await db.storyboards.find_one({"id": storyboard_id}, {"_id": 0})
    if not storyboard:
        raise HTTPException(status_code=404, detail="Storyboard not found")
    
    # Get tiles
    tile_ids = []
    for frame in storyboard.get("frames", []):
        tile_ids.extend(frame.get("tile_refs", []))
    
    tiles = []
    if tile_ids:
        tiles = await db.story_tiles.find({"id": {"$in": tile_ids}}, {"_id": 0}).to_list(100)
    
    return {
        "storyboard": storyboard,
        "tiles": tiles,
        "exported_at": datetime.now(timezone.utc).isoformat()
    }

# Include the router
app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()
